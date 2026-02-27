from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from backend import config  # noqa: E402


def _bool_text(value: bool) -> str:
    return "true" if bool(value) else "false"


def _rgb(r: int, g: int, b: int) -> RGBColor:
    return RGBColor(r, g, b)


def _add_title(slide, text: str, subtitle: str) -> None:
    title = slide.shapes.add_textbox(Inches(0.35), Inches(0.12), Inches(11.5), Inches(0.5))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Segoe UI"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = _rgb(28, 42, 74)

    sub = slide.shapes.add_textbox(Inches(0.35), Inches(0.56), Inches(12.4), Inches(0.3))
    stf = sub.text_frame
    stf.clear()
    sp = stf.paragraphs[0]
    sp.text = subtitle
    sp.font.name = "Segoe UI"
    sp.font.size = Pt(11)
    sp.font.color.rgb = _rgb(84, 98, 125)


def _add_lane(slide, x: float, y: float, w: float, h: float, title: str, fill: RGBColor, border: RGBColor) -> None:
    lane = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    lane.fill.solid()
    lane.fill.fore_color.rgb = fill
    lane.line.color.rgb = border
    lane.line.width = Pt(1.2)

    label = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.04), Inches(3.5), Inches(0.3))
    tf = label.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Segoe UI"
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.color.rgb = _rgb(50, 69, 99)


def _add_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    lines: list[str],
    fill: RGBColor,
    border: RGBColor,
    title_size: int = 12,
    body_size: int = 8,
) -> dict[str, float]:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = border
    box.line.width = Pt(1.1)

    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.07)
    tf.margin_right = Inches(0.07)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.name = "Segoe UI"
    p0.font.size = Pt(title_size)
    p0.font.bold = True
    p0.font.color.rgb = _rgb(31, 45, 76)

    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.name = "Segoe UI"
        p.font.size = Pt(body_size)
        p.font.color.rgb = _rgb(44, 58, 86)
        p.space_after = Pt(0)

    return {"x": x, "y": y, "w": w, "h": h}


def _add_chevron(slide, x: float, y: float, w: float = 0.15, h: float = 0.23) -> None:
    chev = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h))
    chev.fill.solid()
    chev.fill.fore_color.rgb = _rgb(106, 124, 153)
    chev.line.color.rgb = _rgb(106, 124, 153)


def _add_connector(slide, x1: float, y1: float, x2: float, y2: float, text: str) -> None:
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = _rgb(113, 127, 145)
    line.line.width = Pt(1.2)

    label_w = 2.3
    label = slide.shapes.add_textbox(
        Inches((x1 + x2) / 2 - (label_w / 2)),
        Inches((y1 + y2) / 2 - 0.18),
        Inches(label_w),
        Inches(0.35),
    )
    tf = label.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Segoe UI"
    p.font.size = Pt(8)
    p.font.color.rgb = _rgb(79, 92, 112)


def _add_extraction_detail_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(248, 250, 255)
    bg.line.fill.background()

    subtitle = (
        "Detailed ingestion/extraction map from backend/ingestion/* "
        f"(OCR=PaddleOCR, YOLO model={config.YOLO_MODEL})"
    )
    _add_title(slide, "Extraction Step Deep Dive (By Document Type)", subtitle)

    _add_lane(
        slide,
        x=0.25,
        y=0.95,
        w=12.82,
        h=4.05,
        title="Per-Format Extractors (dispatcher in ingest_file())",
        fill=_rgb(236, 243, 255),
        border=_rgb(164, 193, 238),
    )

    _add_box(
        slide,
        x=0.45,
        y=1.22,
        w=12.4,
        h=0.6,
        title="Dispatcher",
        lines=[
            "Routes by extension: .pdf | .pptx | .docx | spreadsheet family | image family | text-like family | generic fallback",
        ],
        fill=_rgb(255, 255, 255),
        border=_rgb(164, 193, 238),
        title_size=10,
        body_size=8,
    )

    card_w = 4.05
    card_h = 1.43
    x1, x2, x3 = 0.45, 4.64, 8.83
    y1, y2 = 1.95, 3.5

    _add_box(
        slide,
        x=x1,
        y=y1,
        w=card_w,
        h=card_h,
        title="PDF (.pdf) -> extract_pdf()",
        lines=[
            "Library: PyMuPDF (fitz)",
            "Text: page.get_text('text') per page",
            "Image render: page.get_pixmap(dpi=PDF_RENDER_DPI)",
            "Sets OCR fallback when native chars < OCR_NATIVE_TEXT_MIN_CHARS",
        ],
        fill=_rgb(255, 255, 255),
        border=_rgb(151, 182, 230),
        title_size=10,
        body_size=7,
    )
    _add_box(
        slide,
        x=x2,
        y=y1,
        w=card_w,
        h=card_h,
        title="PPTX (.pptx) -> extract_pptx()",
        lines=[
            "Library: python-pptx (shapes/text/tables/images)",
            "Slide graph: reading-order + connector edges",
            "Doc graph: slide-to-slide + topic-overlap edges",
            "Optional: LibreOffice (soffice) render -> PDF -> fitz images",
        ],
        fill=_rgb(255, 255, 255),
        border=_rgb(151, 182, 230),
        title_size=10,
        body_size=7,
    )
    _add_box(
        slide,
        x=x3,
        y=y1,
        w=card_w,
        h=card_h,
        title="DOCX (.docx) -> extract_docx()",
        lines=[
            "Primary path: python-docx",
            "Extracts paragraphs, tables, embedded images",
            "Fallback path: OOXML zip + XML parsing",
            "Reads word/document.xml + header/footer XML",
        ],
        fill=_rgb(255, 255, 255),
        border=_rgb(151, 182, 230),
        title_size=10,
        body_size=7,
    )
    _add_box(
        slide,
        x=x1,
        y=y2,
        w=card_w,
        h=card_h,
        title="Spreadsheets (.xlsx/.xlsm/.xltx/.xltm/.xls)",
        lines=[
            "xlsx-family: openpyxl (data_only=True)",
            "xls: xlrd",
            "Rows converted to table text by sheet",
            "openpyxl embedded images extracted when available",
        ],
        fill=_rgb(255, 255, 255),
        border=_rgb(151, 182, 230),
        title_size=10,
        body_size=7,
    )
    _add_box(
        slide,
        x=x2,
        y=y2,
        w=card_w,
        h=card_h,
        title="Images (.png/.jpg/.jpeg/.tiff/.bmp/.gif/.webp)",
        lines=[
            "Loader: PIL Image.open() -> image block",
            "OCR candidate path enabled by metadata flag",
            "Diagram parser always attempted when enabled",
            "Can emit OCR text + graph/node/edge evidence chunks",
        ],
        fill=_rgb(255, 255, 255),
        border=_rgb(151, 182, 230),
        title_size=10,
        body_size=7,
    )
    _add_box(
        slide,
        x=x3,
        y=y2,
        w=card_w,
        h=card_h,
        title="Text-Like + Generic Fallback",
        lines=[
            "Text-like (.txt/.md/.csv/.json/.yaml/.toml...): read_text utf-8",
            "Generic: HTML/XML (BeautifulSoup or regex), RTF (striprtf), ODF (zip+xml)",
            "Zip signature routes Office files to DOCX/XLSX/PPTX extractors",
            "Last fallback: try PDF parse, then plain text",
        ],
        fill=_rgb(255, 255, 255),
        border=_rgb(151, 182, 230),
        title_size=10,
        body_size=7,
    )

    _add_lane(
        slide,
        x=0.25,
        y=5.1,
        w=12.82,
        h=2.1,
        title="Shared Image Enrichment Internals",
        fill=_rgb(235, 249, 239),
        border=_rgb(159, 212, 178),
    )

    b1 = _add_box(
        slide,
        x=0.45,
        y=5.45,
        w=3.95,
        h=1.45,
        title="OCR Path (backend/ingestion/ocr.py + ocr_worker.py)",
        lines=[
            f"Engine: PaddleOCR (lang={config.PADDLE_OCR_LANG}, use_gpu={_bool_text(config.PADDLE_OCR_USE_GPU)})",
            "Controller encodes image as base64 PNG",
            "JSON request/response over stdin/stdout with isolated subprocess",
            "Writes ocr_status/ocr_engine/ocr_confidence/ocr_line_count/ocr_error",
        ],
        fill=_rgb(255, 255, 255),
        border=_rgb(142, 198, 163),
        title_size=9,
        body_size=7,
    )
    b2 = _add_box(
        slide,
        x=4.66,
        y=5.45,
        w=4.02,
        h=1.45,
        title="Diagram Graph Path (backend/ingestion/diagram_parser.py)",
        lines=[
            f"Node detector: Ultralytics YOLO model {config.YOLO_MODEL}; fallback: OpenCV contours",
            "Node labels: OCR per detected node crop",
            "Edge extraction: OpenCV Canny + HoughLinesP, attach nearest nodes",
            "Graph metrics via NetworkX when installed (components, largest, density)",
        ],
        fill=_rgb(255, 255, 255),
        border=_rgb(142, 198, 163),
        title_size=9,
        body_size=7,
    )
    b3 = _add_box(
        slide,
        x=8.93,
        y=5.45,
        w=3.92,
        h=1.45,
        title="Chunk Outputs + Persistence",
        lines=[
            "Image text + OCR merged, then chunk_text(max_chars=900, overlap=120)",
            "Diagram chunks: diagram_graph, diagram_node, diagram_edge",
            "Stores graph JSON in diagram_graphs table and chunks in chunks table",
            "Embeddings created next and persisted to embeddings table",
        ],
        fill=_rgb(255, 255, 255),
        border=_rgb(142, 198, 163),
        title_size=9,
        body_size=7,
    )

    _add_chevron(slide, x=b1["x"] + b1["w"] + 0.09, y=6.03, w=0.18, h=0.27)
    _add_chevron(slide, x=b2["x"] + b2["w"] + 0.09, y=6.03, w=0.18, h=0.27)


def build_presentation(output_path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(247, 249, 253)
    bg.line.fill.background()

    subtitle = (
        "Source: docs/doc_chatbot_workflow.md + live config "
        f"(chat={config.OPENAI_CHAT_MODEL}, embed={config.OPENAI_EMBED_MODEL}, router={config.OPENAI_ROUTER_MODEL})"
    )
    _add_title(slide, "doc_chatbot RAG Workflow (Model-Aware Architecture)", subtitle)

    _add_lane(
        slide,
        x=0.25,
        y=0.9,
        w=12.82,
        h=2.9,
        title="Ingestion Pipeline",
        fill=_rgb(234, 243, 255),
        border=_rgb(163, 196, 236),
    )
    _add_lane(
        slide,
        x=0.25,
        y=3.95,
        w=12.82,
        h=3.25,
        title="Question Answering Pipeline",
        fill=_rgb(234, 248, 239),
        border=_rgb(164, 213, 182),
    )

    ingest_y = 1.32
    ingest_h = 1.95
    ingest_boxes: list[dict[str, float]] = []
    ingest_boxes.append(
        _add_box(
            slide,
            x=0.35,
            y=ingest_y,
            w=1.55,
            h=ingest_h,
            title="1) Upload/API",
            lines=[
                "POST /api/documents",
                "FastAPI background task",
                "queued -> processing -> ready",
            ],
            fill=_rgb(255, 255, 255),
            border=_rgb(146, 182, 226),
            body_size=8,
        )
    )
    ingest_boxes.append(
        _add_box(
            slide,
            x=2.15,
            y=ingest_y,
            w=2.45,
            h=ingest_h,
            title="2) Extract + Enrich",
            lines=[
                "Parsers: PyMuPDF, python-pptx,",
                "python-docx, openpyxl/xlrd",
                f"OCR: PaddleOCR (lang={config.PADDLE_OCR_LANG}, GPU={_bool_text(config.PADDLE_OCR_USE_GPU)})",
                f"Nodes: Ultralytics YOLO ({config.YOLO_MODEL})",
                "Edges: OpenCV Canny + Hough",
            ],
            fill=_rgb(255, 255, 255),
            border=_rgb(146, 182, 226),
            body_size=8,
        )
    )
    ingest_boxes.append(
        _add_box(
            slide,
            x=4.85,
            y=ingest_y,
            w=1.45,
            h=ingest_h,
            title="3) Chunking",
            lines=[
                "chunk_text()",
                "about 900 chars",
                "120-char overlap",
            ],
            fill=_rgb(255, 255, 255),
            border=_rgb(146, 182, 226),
            body_size=8,
        )
    )
    ingest_boxes.append(
        _add_box(
            slide,
            x=6.55,
            y=ingest_y,
            w=1.8,
            h=ingest_h,
            title="4) Embeddings",
            lines=[
                "Provider: OpenAI",
                f"Model: {config.OPENAI_EMBED_MODEL}",
                "normalized vectors",
                f"dim={config.PGVECTOR_DIM}",
            ],
            fill=_rgb(255, 255, 255),
            border=_rgb(146, 182, 226),
            body_size=8,
        )
    )
    ingest_boxes.append(
        _add_box(
            slide,
            x=8.6,
            y=ingest_y,
            w=2.0,
            h=ingest_h,
            title="5) Persistence",
            lines=[
                "documents / chunks /",
                "embeddings / diagram_graphs",
                f"DB backend: {config.DB_BACKEND}",
                "pgvector in Postgres mode",
            ],
            fill=_rgb(255, 255, 255),
            border=_rgb(146, 182, 226),
            body_size=8,
        )
    )
    ingest_boxes.append(
        _add_box(
            slide,
            x=10.85,
            y=ingest_y,
            w=1.95,
            h=ingest_h,
            title="6) Indexes",
            lines=[
                "Dense: pgvector SQL or",
                "in-memory dot product",
                "Sparse: BM25 regex",
                "ready for hybrid search",
            ],
            fill=_rgb(255, 255, 255),
            border=_rgb(146, 182, 226),
            body_size=8,
        )
    )

    ingest_arrow_y = ingest_y + (ingest_h / 2) - 0.11
    for idx in range(len(ingest_boxes) - 1):
        left = ingest_boxes[idx]
        _add_chevron(slide, x=left["x"] + left["w"] + 0.05, y=ingest_arrow_y)

    query_y = 4.5
    query_h = 2.15
    query_boxes: list[dict[str, float]] = []
    query_boxes.append(
        _add_box(
            slide,
            x=0.35,
            y=query_y,
            w=1.35,
            h=query_h,
            title="A) User Query",
            lines=[
                "message +",
                "selected doc_ids",
            ],
            fill=_rgb(255, 255, 255),
            border=_rgb(148, 205, 166),
            body_size=9,
        )
    )
    query_boxes.append(
        _add_box(
            slide,
            x=1.95,
            y=query_y,
            w=1.7,
            h=query_h,
            title="B) Router",
            lines=[
                "OpenAI router",
                f"Model: {config.OPENAI_ROUTER_MODEL}",
                "task + strategy + top_k",
            ],
            fill=_rgb(255, 255, 255),
            border=_rgb(148, 205, 166),
            body_size=8,
        )
    )
    query_boxes.append(
        _add_box(
            slide,
            x=3.9,
            y=query_y,
            w=2.2,
            h=query_h,
            title="C) Retrieval",
            lines=[
                f"Mode: {config.RETRIEVAL_MODE}",
                "Dense + BM25 + RRF fusion",
                f"TOP_K={config.TOP_K}",
                "doc scope validation",
            ],
            fill=_rgb(255, 255, 255),
            border=_rgb(148, 205, 166),
            body_size=8,
        )
    )
    query_boxes.append(
        _add_box(
            slide,
            x=6.35,
            y=query_y,
            w=2.1,
            h=query_h,
            title="D) Reranking",
            lines=[
                f"CrossEncoder: {config.RERANK_MODEL_ID}",
                f"enabled={_bool_text(config.ENABLE_RERANKER)}",
                "fallback: lexical overlap",
            ],
            fill=_rgb(255, 255, 255),
            border=_rgb(148, 205, 166),
            body_size=8,
        )
    )
    query_boxes.append(
        _add_box(
            slide,
            x=8.7,
            y=query_y,
            w=2.0,
            h=query_h,
            title="E) Context Build",
            lines=[
                "diagram-aware evidence mix",
                "source-tagged blocks",
                f"MAX_CONTEXT_CHARS={config.MAX_CONTEXT_CHARS}",
            ],
            fill=_rgb(255, 255, 255),
            border=_rgb(148, 205, 166),
            body_size=8,
        )
    )
    query_boxes.append(
        _add_box(
            slide,
            x=10.95,
            y=query_y,
            w=1.85,
            h=query_h,
            title="F) Answer",
            lines=[
                f"OpenAI chat: {config.OPENAI_CHAT_MODEL}",
                "multimodal context",
                "grounded response + sources",
            ],
            fill=_rgb(255, 255, 255),
            border=_rgb(148, 205, 166),
            body_size=8,
        )
    )

    query_arrow_y = query_y + (query_h / 2) - 0.11
    for idx in range(len(query_boxes) - 1):
        left = query_boxes[idx]
        _add_chevron(slide, x=left["x"] + left["w"] + 0.05, y=query_arrow_y)

    _add_connector(
        slide,
        x1=11.85,
        y1=3.25,
        x2=5.0,
        y2=4.5,
        text="retrieval reads chunk + index state",
    )

    note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.0), Inches(0.97), Inches(4.0), Inches(0.9))
    note.fill.solid()
    note.fill.fore_color.rgb = _rgb(255, 247, 222)
    note.line.color.rgb = _rgb(224, 191, 97)
    note.line.width = Pt(1.0)
    tf = note.text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = "Runtime Model Notes"
    p0.font.name = "Segoe UI"
    p0.font.bold = True
    p0.font.size = Pt(10)
    p0.font.color.rgb = _rgb(117, 82, 4)
    for line in [
        f"OCR enabled={_bool_text(config.ENABLE_OCR)}",
        f"YOLO enabled={_bool_text(config.ENABLE_YOLO_DIAGRAM_DETECTOR)}",
        f"AI ingest summaries enabled={_bool_text(config.ENABLE_AI_INGEST_SUMMARIES)} (if true -> gpt-4o-mini)",
    ]:
        p = tf.add_paragraph()
        p.text = line
        p.font.name = "Segoe UI"
        p.font.size = Pt(8)
        p.font.color.rgb = _rgb(117, 82, 4)
        p.space_after = Pt(0)

    _add_extraction_detail_slide(prs)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def main() -> None:
    output = build_presentation(ROOT / "docs" / "doc_chatbot_architecture_diagram.pptx")
    print(output)


if __name__ == "__main__":
    main()
