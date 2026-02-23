from __future__ import annotations

import io
import logging
import math
import re
import shutil
import subprocess
import tempfile
import zipfile
from collections import Counter
from pathlib import Path
from typing import Dict, List, Tuple
from xml.etree import ElementTree as ET

import fitz
from PIL import Image
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation

from .. import config

try:
    from bs4 import BeautifulSoup
except Exception:  # pragma: no cover - optional dependency
    BeautifulSoup = None

try:
    from docx import Document
except Exception:  # pragma: no cover - optional dependency
    Document = None

try:
    from openpyxl import load_workbook
except Exception:  # pragma: no cover - optional dependency
    load_workbook = None

try:
    import xlrd
except Exception:  # pragma: no cover - optional dependency
    xlrd = None

try:
    from striprtf.striprtf import rtf_to_text
except Exception:  # pragma: no cover - optional dependency
    rtf_to_text = None


Block = Dict[str, object]
DOCX_NS = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
ODF_TEXT_NS = {
    "text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0",
    "table": "urn:oasis:names:tc:opendocument:xmlns:table:1.0",
}
LOGGER = logging.getLogger(__name__)
PPTX_STOPWORDS = {
    "the",
    "and",
    "for",
    "with",
    "that",
    "this",
    "from",
    "into",
    "have",
    "has",
    "are",
    "was",
    "were",
    "your",
    "their",
    "slide",
    "slides",
    "about",
    "there",
    "they",
    "them",
    "what",
    "when",
    "which",
    "will",
    "would",
    "could",
    "should",
}


def _table_to_text(rows: List[List[object]]) -> str:
    cleaned_rows: List[List[str]] = []
    width = 0
    for row in rows:
        values = [str(cell).strip() if cell is not None else "" for cell in row]
        width = max(width, len(values))
        cleaned_rows.append(values)
    if width == 0:
        return ""
    normalized = [row + ([""] * (width - len(row))) for row in cleaned_rows]
    table_lines = ["\t".join(row).rstrip() for row in normalized if any(cell for cell in row)]
    return "\n".join(table_lines).strip()


def _image_from_blob(image_blob: bytes) -> Image.Image | None:
    try:
        image = Image.open(io.BytesIO(image_blob))
        image.load()
        return image
    except Exception:
        return None


def _normalized_bbox(shape, slide_width: int, slide_height: int) -> Dict[str, float]:
    width = max(1, int(slide_width or 1))
    height = max(1, int(slide_height or 1))
    left = int(getattr(shape, "left", 0) or 0)
    top = int(getattr(shape, "top", 0) or 0)
    shape_width = max(0, int(getattr(shape, "width", 0) or 0))
    shape_height = max(0, int(getattr(shape, "height", 0) or 0))
    return {
        "x": round(left / float(width), 5),
        "y": round(top / float(height), 5),
        "w": round(shape_width / float(width), 5),
        "h": round(shape_height / float(height), 5),
        "left_emu": left,
        "top_emu": top,
        "width_emu": shape_width,
        "height_emu": shape_height,
    }


def _shape_center_emu(shape) -> Tuple[int, int]:
    left = int(getattr(shape, "left", 0) or 0)
    top = int(getattr(shape, "top", 0) or 0)
    width = int(getattr(shape, "width", 0) or 0)
    height = int(getattr(shape, "height", 0) or 0)
    return (left + (width // 2), top + (height // 2))


def _find_soffice_binary() -> str | None:
    candidates = [
        "soffice",
        "libreoffice",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]
    for candidate in candidates:
        resolved = shutil.which(candidate)
        if resolved:
            return resolved
        candidate_path = Path(candidate)
        if candidate_path.exists():
            return str(candidate_path)
    return None


def _render_pptx_slides(path: Path) -> List[Block]:
    if not config.ENABLE_PPTX_SLIDE_RENDER:
        return []
    soffice = _find_soffice_binary()
    if not soffice:
        LOGGER.info("Skipping PPTX full-slide rendering because LibreOffice/soffice was not found on PATH.")
        return []

    blocks: List[Block] = []
    with tempfile.TemporaryDirectory(prefix="pptx_render_") as tmp_dir:
        target_dir = Path(tmp_dir)
        command = [
            soffice,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(target_dir),
            str(path),
        ]
        try:
            result = subprocess.run(
                command,
                check=False,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
            )
        except Exception as exc:
            LOGGER.warning("Failed to render PPTX with LibreOffice: %s", exc)
            return blocks
        if result.returncode != 0:
            stderr = (result.stderr or "").strip()
            LOGGER.warning("LibreOffice conversion failed for %s: %s", path.name, stderr or "unknown error")
            return blocks

        output_pdf = target_dir / f"{path.stem}.pdf"
        if not output_pdf.exists():
            candidates = sorted(target_dir.glob("*.pdf"))
            if not candidates:
                return blocks
            output_pdf = candidates[0]

        try:
            with fitz.open(str(output_pdf)) as rendered:
                for slide_index, page in enumerate(rendered, start=1):
                    pix = page.get_pixmap(dpi=config.PDF_RENDER_DPI)
                    if pix.alpha:
                        pix = fitz.Pixmap(fitz.csRGB, pix)
                    image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    blocks.append(
                        {
                            "page": slide_index,
                            "type": "image",
                            "image": image,
                            "metadata": {
                                "image_kind": "slide_render",
                                "ocr_fallback_required": True,
                                "native_text_chars": 0,
                                "render_backend": "libreoffice_pdf",
                            },
                        }
                    )
        except Exception as exc:
            LOGGER.warning("Failed to process rendered PPTX PDF for %s: %s", path.name, exc)

    return blocks


def _extract_keywords(text: str, limit: int = 8) -> List[str]:
    tokens = re.findall(r"[A-Za-z0-9_]{3,}", str(text or "").lower())
    counts = Counter(token for token in tokens if token not in PPTX_STOPWORDS)
    return [token for token, _ in counts.most_common(limit)]


def _nearest_node(point: Tuple[int, int], nodes: List[Dict[str, object]]) -> Dict[str, object] | None:
    if not nodes:
        return None
    px, py = point
    best_node = None
    best_dist = None
    for node in nodes:
        nx = int(node.get("center_x_emu") or 0)
        ny = int(node.get("center_y_emu") or 0)
        dist = ((px - nx) ** 2) + ((py - ny) ** 2)
        if best_dist is None or dist < best_dist:
            best_dist = dist
            best_node = node
    return best_node


def _slide_graph_to_text(graph: Dict[str, object]) -> str:
    slide_index = graph.get("slide_index")
    edges = graph.get("edges")
    nodes = graph.get("nodes")
    relationships: List[str] = []
    if isinstance(edges, list):
        for edge in edges[:20]:
            if not isinstance(edge, dict):
                continue
            edge_type = str(edge.get("type") or "relation")
            source = str(edge.get("from") or "?")
            target = str(edge.get("to") or "?")
            relationships.append(f"{source} -> {target} ({edge_type})")
    key_nodes: List[str] = []
    if isinstance(nodes, list):
        for node in nodes[:15]:
            if not isinstance(node, dict):
                continue
            label = str(node.get("label") or node.get("id") or "").strip()
            node_kind = str(node.get("node_kind") or "shape")
            if label:
                key_nodes.append(f"{label} ({node_kind})")
    lines = [
        f"Slide {slide_index} relationship graph.",
        f"Node count: {graph.get('node_count', 0)}. Edge count: {graph.get('edge_count', 0)}.",
    ]
    if key_nodes:
        lines.append("Key nodes: " + "; ".join(key_nodes))
    if relationships:
        lines.append("Relationships: " + "; ".join(relationships))
    return "\n".join(lines).strip()


def _document_graph_to_text(graph: Dict[str, object]) -> str:
    slide_nodes = graph.get("slide_nodes")
    edges = graph.get("edges")
    lines = ["Document slide relationship graph."]
    if isinstance(slide_nodes, list):
        summary = []
        for node in slide_nodes[:20]:
            if not isinstance(node, dict):
                continue
            idx = node.get("slide_index")
            keywords = node.get("keywords") if isinstance(node.get("keywords"), list) else []
            summary.append(f"Slide {idx}: {', '.join(str(k) for k in keywords[:6])}")
        if summary:
            lines.append("Slide topics: " + "; ".join(summary))
    if isinstance(edges, list):
        edge_summary = []
        for edge in edges[:40]:
            if not isinstance(edge, dict):
                continue
            source = edge.get("from")
            target = edge.get("to")
            edge_type = edge.get("type")
            if source and target:
                edge_summary.append(f"{source}->{target} ({edge_type})")
        if edge_summary:
            lines.append("Slide-to-slide relationships: " + "; ".join(edge_summary))
    return "\n".join(lines).strip()


def extract_pdf(path: Path) -> List[Block]:
    blocks: List[Block] = []

    with fitz.open(str(path)) as doc:
        for page_index, page in enumerate(doc, start=1):
            native_text = (page.get_text("text") or "").strip()
            native_non_ws_chars = len(re.sub(r"\s+", "", native_text))
            needs_ocr_fallback = native_non_ws_chars < config.OCR_NATIVE_TEXT_MIN_CHARS

            if native_text:
                blocks.append(
                    {
                        "page": page_index,
                        "type": "text",
                        "text": native_text,
                        "metadata": {
                            "extraction_method": "native_pdf_text",
                            "native_text_chars": native_non_ws_chars,
                        },
                    }
                )

            # Keep page images for visual follow-up and OCR fallback on scanned pages.
            pix = page.get_pixmap(dpi=config.PDF_RENDER_DPI)
            if pix.alpha:
                pix = fitz.Pixmap(fitz.csRGB, pix)
            image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            blocks.append(
                {
                    "page": page_index,
                    "type": "image",
                    "image": image,
                    "metadata": {
                        "ocr_fallback_required": needs_ocr_fallback,
                        "native_text_chars": native_non_ws_chars,
                    },
                }
            )

    return blocks


def extract_pptx(path: Path) -> List[Block]:
    blocks: List[Block] = []
    prs = Presentation(str(path))
    slide_width = max(1, int(prs.slide_width or 1))
    slide_height = max(1, int(prs.slide_height or 1))
    doc_slide_nodes: List[Dict[str, object]] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_text_parts: List[str] = []
        graph_nodes: List[Dict[str, object]] = []
        graph_edges: List[Dict[str, object]] = []
        connectors: List[Dict[str, object]] = []

        for shape in slide.shapes:
            shape_id = int(getattr(shape, "shape_id", 0) or 0)
            node_id = f"slide:{slide_index}:shape:{shape_id or (len(graph_nodes) + 1)}"
            shape_name = str(getattr(shape, "name", node_id) or node_id)
            shape_type_enum = getattr(shape, "shape_type", None)
            shape_type_name = str(getattr(shape_type_enum, "name", shape_type_enum or "unknown"))
            bbox = _normalized_bbox(shape, slide_width, slide_height)
            center_x, center_y = _shape_center_emu(shape)

            node: Dict[str, object] = {
                "id": node_id,
                "shape_id": shape_id,
                "label": shape_name,
                "shape_type": shape_type_name,
                "node_kind": "shape",
                "bbox": bbox,
                "center_x_emu": center_x,
                "center_y_emu": center_y,
            }

            common_metadata = {
                "slide_index": slide_index,
                "shape_id": shape_id,
                "shape_name": shape_name,
                "shape_type": shape_type_name,
                "bbox": bbox,
            }

            text_value = ""
            if hasattr(shape, "text"):
                try:
                    text_value = str(shape.text or "").strip()
                except Exception:
                    text_value = ""
            if text_value:
                node["node_kind"] = "text"
                node["text_excerpt"] = text_value[:240]
                slide_text_parts.append(text_value)
                blocks.append(
                    {
                        "page": slide_index,
                        "type": "text",
                        "text": text_value,
                        "metadata": common_metadata,
                    }
                )

            if getattr(shape, "has_table", False):
                rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                table_text = _table_to_text(rows)
                if table_text:
                    node["node_kind"] = "table"
                    node["text_excerpt"] = table_text[:240]
                    slide_text_parts.append(table_text)
                    blocks.append(
                        {
                            "page": slide_index,
                            "type": "table",
                            "text": table_text,
                            "metadata": common_metadata,
                        }
                    )

            if shape_type_enum == MSO_SHAPE_TYPE.PICTURE:
                image_bytes = shape.image.blob
                image = _image_from_blob(image_bytes)
                if image is not None:
                    node["node_kind"] = "image"
                    blocks.append(
                        {
                            "page": slide_index,
                            "type": "image",
                            "image": image,
                            "metadata": {
                                **common_metadata,
                                "image_kind": "embedded_picture",
                                "ocr_fallback_required": True,
                            },
                        }
                    )

            if shape_type_enum == MSO_SHAPE_TYPE.LINE:
                node["node_kind"] = "connector"
                begin_x = int(getattr(shape, "begin_x", center_x) or center_x)
                begin_y = int(getattr(shape, "begin_y", center_y) or center_y)
                end_x = int(getattr(shape, "end_x", center_x) or center_x)
                end_y = int(getattr(shape, "end_y", center_y) or center_y)
                connectors.append(
                    {
                        "id": node_id,
                        "begin_x": begin_x,
                        "begin_y": begin_y,
                        "end_x": end_x,
                        "end_y": end_y,
                    }
                )

            graph_nodes.append(node)

        text_nodes = [node for node in graph_nodes if str(node.get("text_excerpt") or "").strip()]
        text_nodes = sorted(
            text_nodes,
            key=lambda node: (
                float(((node.get("bbox") or {}).get("y") if isinstance(node.get("bbox"), dict) else 0.0) or 0.0),
                float(((node.get("bbox") or {}).get("x") if isinstance(node.get("bbox"), dict) else 0.0) or 0.0),
            ),
        )
        for idx in range(len(text_nodes) - 1):
            source = text_nodes[idx]
            target = text_nodes[idx + 1]
            graph_edges.append({"type": "reading_order", "from": source.get("id"), "to": target.get("id")})

        non_connector_nodes = [node for node in graph_nodes if node.get("node_kind") != "connector"]
        for connector in connectors:
            start = _nearest_node((connector["begin_x"], connector["begin_y"]), non_connector_nodes)
            end = _nearest_node((connector["end_x"], connector["end_y"]), non_connector_nodes)
            if start is None or end is None:
                continue
            start_id = str(start.get("id") or "")
            end_id = str(end.get("id") or "")
            if not start_id or not end_id or start_id == end_id:
                continue
            distance = math.sqrt(
                ((int(start.get("center_x_emu") or 0) - int(end.get("center_x_emu") or 0)) ** 2)
                + ((int(start.get("center_y_emu") or 0) - int(end.get("center_y_emu") or 0)) ** 2)
            )
            graph_edges.append(
                {
                    "type": "connector",
                    "from": start_id,
                    "to": end_id,
                    "via": connector["id"],
                    "distance_emu": int(distance),
                }
            )

        slide_text = "\n".join(part for part in slide_text_parts if part).strip()
        keywords = _extract_keywords(slide_text)
        title = ""
        if slide_text:
            first_line = slide_text.splitlines()[0].strip()
            title = first_line[:120]

        doc_slide_nodes.append(
            {
                "id": f"slide:{slide_index}",
                "slide_index": slide_index,
                "title": title or f"Slide {slide_index}",
                "keywords": keywords,
            }
        )

        if config.ENABLE_PPTX_RELATIONSHIP_GRAPH:
            max_edges = max(0, int(config.PPTX_GRAPH_MAX_EDGES))
            trimmed_edges = graph_edges[:max_edges] if max_edges else graph_edges
            slide_graph: Dict[str, object] = {
                "kind": "pptx_slide_graph",
                "slide_index": slide_index,
                "node_count": len(graph_nodes),
                "edge_count": len(trimmed_edges),
                "keywords": keywords,
                "nodes": graph_nodes,
                "edges": trimmed_edges,
            }
            blocks.append(
                {
                    "page": slide_index,
                    "type": "slide_graph",
                    "text": _slide_graph_to_text(slide_graph),
                    "graph": slide_graph,
                    "parser_version": "pptx-slide-graph-v1",
                    "metadata": {
                        "graph_scope": "slide",
                        "graph_kind": "pptx_slide_graph",
                        "slide_index": slide_index,
                    },
                }
            )

    if config.ENABLE_PPTX_RELATIONSHIP_GRAPH and doc_slide_nodes:
        max_edges = max(0, int(config.PPTX_GRAPH_MAX_EDGES))
        doc_edges: List[Dict[str, object]] = []
        for idx in range(len(doc_slide_nodes) - 1):
            current_slide = doc_slide_nodes[idx]
            next_slide = doc_slide_nodes[idx + 1]
            doc_edges.append(
                {
                    "type": "next_slide",
                    "from": current_slide.get("id"),
                    "to": next_slide.get("id"),
                }
            )

        for i in range(len(doc_slide_nodes)):
            keywords_i = set(str(k) for k in (doc_slide_nodes[i].get("keywords") or []))
            if not keywords_i:
                continue
            for j in range(i + 1, len(doc_slide_nodes)):
                keywords_j = set(str(k) for k in (doc_slide_nodes[j].get("keywords") or []))
                overlap = sorted(keywords_i.intersection(keywords_j))
                if len(overlap) < 2:
                    continue
                doc_edges.append(
                    {
                        "type": "topic_overlap",
                        "from": doc_slide_nodes[i].get("id"),
                        "to": doc_slide_nodes[j].get("id"),
                        "shared_keywords": overlap[:6],
                        "shared_count": len(overlap),
                    }
                )
                if max_edges and len(doc_edges) >= max_edges:
                    break
            if max_edges and len(doc_edges) >= max_edges:
                break

        if max_edges:
            doc_edges = doc_edges[:max_edges]
        doc_graph: Dict[str, object] = {
            "kind": "pptx_document_graph",
            "slide_count": len(doc_slide_nodes),
            "slide_nodes": doc_slide_nodes,
            "edges": doc_edges,
        }
        blocks.append(
            {
                "page": 1,
                "type": "slide_graph",
                "text": _document_graph_to_text(doc_graph),
                "graph": doc_graph,
                "parser_version": "pptx-document-graph-v1",
                "metadata": {
                    "graph_scope": "document",
                    "graph_kind": "pptx_document_graph",
                },
            }
        )

    rendered_blocks = _render_pptx_slides(path)
    if rendered_blocks:
        blocks.extend(rendered_blocks)

    return blocks


def _docx_text_from_xml(xml_bytes: bytes) -> str:
    root = ET.fromstring(xml_bytes)
    paragraphs: List[str] = []
    for paragraph in root.findall(".//w:p", DOCX_NS):
        pieces: List[str] = []
        for node in paragraph.iter():
            tag = node.tag.rsplit("}", 1)[-1]
            if tag == "t" and node.text:
                pieces.append(node.text)
            elif tag == "tab":
                pieces.append("\t")
            elif tag in {"br", "cr"}:
                pieces.append("\n")
        text = "".join(pieces).strip()
        if text:
            paragraphs.append(text)
    return "\n".join(paragraphs).strip()


def extract_docx(path: Path) -> List[Block]:
    if Document is not None:
        return _extract_docx_with_python_docx(path)

    blocks: List[Block] = []
    with zipfile.ZipFile(path) as archive:
        names = set(archive.namelist())
        part_names = ["word/document.xml"]
        part_names.extend(
            sorted(
                name
                for name in names
                if name.startswith("word/header") and name.endswith(".xml")
            )
        )
        part_names.extend(
            sorted(
                name
                for name in names
                if name.startswith("word/footer") and name.endswith(".xml")
            )
        )

        page = 1
        for part in part_names:
            if part not in names:
                continue
            text = _docx_text_from_xml(archive.read(part))
            if text:
                blocks.append({"page": page, "type": "text", "text": text})
                page += 1
    return blocks


def _extract_docx_with_python_docx(path: Path) -> List[Block]:
    blocks: List[Block] = []
    doc = Document(str(path))

    paragraph_text = "\n".join(p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()).strip()
    if paragraph_text:
        blocks.append({"page": 1, "type": "text", "text": paragraph_text})

    for table_idx, table in enumerate(doc.tables, start=1):
        rows: List[List[str]] = []
        for row in table.rows:
            rows.append([cell.text.strip() for cell in row.cells])
        table_text = _table_to_text(rows)
        if table_text:
            blocks.append({"page": 1, "type": "table", "text": f"Table {table_idx}\n{table_text}"})

    seen_parts = set()
    for rel in doc.part.rels.values():
        if "image" not in rel.reltype:
            continue
        target_part = getattr(rel, "target_part", None)
        partname = str(getattr(target_part, "partname", ""))
        if not target_part or partname in seen_parts:
            continue
        seen_parts.add(partname)
        image = _image_from_blob(target_part.blob)
        if image is not None:
            blocks.append({"page": 1, "type": "image", "image": image})
    return blocks


def extract_xlsx(path: Path) -> List[Block]:
    if load_workbook is None:
        raise RuntimeError("openpyxl is required to parse spreadsheet files")

    blocks: List[Block] = []
    workbook = load_workbook(str(path), data_only=True)
    try:
        for sheet_index, sheet in enumerate(workbook.worksheets, start=1):
            rows: List[List[object]] = []
            for row in sheet.iter_rows(values_only=True):
                values = list(row)
                if any(value not in (None, "") for value in values):
                    rows.append(values)

            table_text = _table_to_text(rows)
            if table_text:
                blocks.append(
                    {
                        "page": sheet_index,
                        "type": "table",
                        "text": f"Sheet: {sheet.title}\n{table_text}",
                    }
                )

            for image_ref in getattr(sheet, "_images", []):
                image_bytes = None
                if hasattr(image_ref, "_data"):
                    try:
                        image_bytes = image_ref._data()
                    except Exception:
                        image_bytes = None
                if not image_bytes:
                    continue
                image = _image_from_blob(image_bytes)
                if image is not None:
                    blocks.append({"page": sheet_index, "type": "image", "image": image})
    finally:
        workbook.close()
    return blocks


def extract_xls(path: Path) -> List[Block]:
    if xlrd is None:
        raise RuntimeError("xlrd is required to parse .xls spreadsheet files")

    blocks: List[Block] = []
    workbook = xlrd.open_workbook(str(path), on_demand=True)
    try:
        for sheet_index in range(workbook.nsheets):
            sheet = workbook.sheet_by_index(sheet_index)
            rows: List[List[object]] = []
            for row_index in range(sheet.nrows):
                values = sheet.row_values(row_index)
                if any(value not in (None, "") for value in values):
                    rows.append(values)
            table_text = _table_to_text(rows)
            if table_text:
                blocks.append(
                    {
                        "page": sheet_index + 1,
                        "type": "table",
                        "text": f"Sheet: {sheet.name}\n{table_text}",
                    }
                )
    finally:
        workbook.release_resources()
    return blocks


def extract_html(path: Path) -> List[Block]:
    raw = path.read_text(encoding="utf-8", errors="ignore")
    if BeautifulSoup is not None:
        soup = BeautifulSoup(raw, "html.parser")
        text = soup.get_text(separator="\n").strip()
    else:
        text = re.sub(r"<[^>]+>", " ", raw)
        text = re.sub(r"\s+", " ", text).strip()
    if not text:
        return []
    return [{"page": 1, "type": "text", "text": text}]


def extract_rtf(path: Path) -> List[Block]:
    raw = path.read_text(encoding="utf-8", errors="ignore")
    text = rtf_to_text(raw).strip() if rtf_to_text is not None else raw.strip()
    if not text:
        return []
    return [{"page": 1, "type": "text", "text": text}]


def extract_odf(path: Path) -> List[Block]:
    blocks: List[Block] = []
    with zipfile.ZipFile(path) as archive:
        if "content.xml" not in archive.namelist():
            return blocks
        content_xml = archive.read("content.xml")
    root = ET.fromstring(content_xml)
    text_parts = [
        node.text.strip()
        for node in root.findall(".//text:p", ODF_TEXT_NS)
        if node.text and node.text.strip()
    ]
    if text_parts:
        blocks.append({"page": 1, "type": "text", "text": "\n".join(text_parts)})

    table_rows: List[List[object]] = []
    for row_node in root.findall(".//table:table-row", ODF_TEXT_NS):
        row_values: List[str] = []
        for cell_node in row_node.findall(".//table:table-cell", ODF_TEXT_NS):
            cell_text = "".join(cell_node.itertext()).strip()
            row_values.append(cell_text)
        if any(value for value in row_values):
            table_rows.append(row_values)
    table_text = _table_to_text(table_rows)
    if table_text:
        blocks.append({"page": 1, "type": "table", "text": table_text})
    return blocks


def extract_image(path: Path) -> List[Block]:
    image = Image.open(path)
    blocks: List[Block] = []
    blocks.append({"page": 1, "type": "image", "image": image})
    return blocks


def extract_text(path: Path) -> List[Block]:
    text = path.read_text(encoding="utf-8", errors="ignore").strip()
    if not text:
        return []
    return [{"page": 1, "type": "text", "text": text}]


def extract_generic(path: Path) -> List[Block]:
    suffix = path.suffix.lower()
    if suffix in {".htm", ".html", ".xhtml", ".xml"}:
        return extract_html(path)
    if suffix in {".rtf"}:
        return extract_rtf(path)
    if suffix in {".xlsx", ".xlsm", ".xltx", ".xltm"}:
        return extract_xlsx(path)
    if suffix in {".xls"}:
        return extract_xls(path)
    if suffix in {".odt", ".ods", ".odp"}:
        return extract_odf(path)

    if zipfile.is_zipfile(path):
        with zipfile.ZipFile(path) as archive:
            names = set(archive.namelist())
            if "word/document.xml" in names:
                return extract_docx(path)
            if "xl/workbook.xml" in names:
                return extract_xlsx(path)
            if "ppt/presentation.xml" in names:
                return extract_pptx(path)
            if "content.xml" in names:
                return extract_odf(path)

    try:
        return extract_pdf(path)
    except Exception:
        pass

    blocks = extract_text(path)
    if blocks:
        return blocks

    ext = suffix or "[no extension]"
    raise ValueError(f"Unsupported file type: {ext}")

